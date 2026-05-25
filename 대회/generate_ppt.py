"""
오늘시간표 - 제8회 교육 공공데이터 AI활용대회 PPT 생성
python3 generate_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === 색상 ===
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1A, 0x25, 0x3C)
ACCENT = RGBColor(0x38, 0x7A, 0xFF)
ACCENT2 = RGBColor(0x00, 0xC8, 0x9D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9E, 0xA3, 0xB5)
LIGHT = RGBColor(0xE8, 0xEB, 0xF0)
ORANGE = RGBColor(0xFF, 0x8C, 0x42)
RED = RGBColor(0xFF, 0x5C, 0x5C)
PURPLE = RGBColor(0xA8, 0x78, 0xFF)
GREEN = RGBColor(0x34, 0xC7, 0x59)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, color=BG_CARD, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_text(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_para(tf, text, size=14, color=LIGHT, bold=False, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(space_before)
    return p

def add_badge(slide, left, top, text, color=ACCENT):
    w, h = Inches(1.8), Inches(0.35)
    shape = add_shape(slide, left, top, w, h, color, 0.5)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.word_wrap = False

# ══════════════════════════════════════
# 1. 표지
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
         "오늘시간표", 54, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(0.6),
         "AI 기반 학교생활 통합 도우미", 24, ACCENT, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.5),
         "교육 공공데이터(NEIS/1365/학교알리미) + AI(Gemini/Groq) + 멀티플랫폼", 14, GRAY, False, PP_ALIGN.CENTER)

# 하단 배지들
badges = ["iOS", "macOS", "watchOS", "위젯", "Live Activity"]
for i, b in enumerate(badges):
    add_badge(slide, Inches(2.8 + i * 1.6), Inches(4.5), b, ACCENT if i < 3 else ACCENT2)

add_text(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.4),
         "제8회 교육 공공데이터 AI활용대회 | 학생부", 14, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.4),
         "팀명: (팀명 입력)", 16, LIGHT, True, PP_ALIGN.CENTER)

# ══════════════════════════════════════
# 2. 활용 데이터 (필수)
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "01  활용 교육 공공데이터", 28, WHITE, True)

data_items = [
    ("NEIS 학교기본정보", "교육부", "학교 검색, 유형/주소 확인"),
    ("NEIS 시간표 (초/중/고)", "교육부", "시간표 조회"),
    ("NEIS 급식식단정보", "교육부", "급식 메뉴/칼로리/알레르기"),
    ("NEIS 학사일정", "교육부", "학사일정, 시험 D-Day"),
    ("1365 봉사활동 실적", "행정안전부", "봉사 검색/상세/AI 추천"),
    ("학교알리미 정보공시 (12종)", "교육부", "학생수/교원/동아리/도서관/교복/방과후 등"),
    ("학교알리미 정보공시 (12종)", "교육부", "학생수/교원/동아리/도서관 등"),
]

for i, (name, org, desc) in enumerate(data_items):
    y = Inches(1.3 + i * 0.75)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.65), BG_CARD)
    add_text(slide, Inches(1.0), y + Inches(0.08), Inches(4), Inches(0.25),
             name, 14, WHITE, True)
    add_text(slide, Inches(1.0), y + Inches(0.33), Inches(4), Inches(0.25),
             desc, 11, GRAY)
    add_text(slide, Inches(9.5), y + Inches(0.15), Inches(2.5), Inches(0.3),
             org, 12, ACCENT, True, PP_ALIGN.RIGHT)

add_text(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.3),
         "라이선스: 공공누리 제1유형 (출처 표시)  |  총 7종 교육 공공데이터 활용", 11, GRAY)

# ══════════════════════════════════════
# 3. 문제 인식
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "02  문제 인식", 28, WHITE, True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
         "학생들이 매일 겪는 학교생활의 불편함", 16, GRAY)

problems = [
    ("시간표 확인이 번거롭다", "교실 게시판, 학교 홈페이지, 카톡... 매번 다른 곳에서 확인해야 함", "calendar"),
    ("급식/학사일정 정보가 분산", "학교 홈페이지 접속 필요, 알레르기 확인 불편, D-Day 수동 계산", "fork.knife"),
    ("봉사활동 관리가 어렵다", "1365 웹사이트 검색 복잡, 지역/관심 분야 필터 부족", "hand.raised"),
    ("수행평가 준비에 도움이 없다", "종이 공지 의존, 준비 방법/평가기준 파악 어려움", "pencil"),
    ("학교 정보를 한눈에 볼 수 없다", "학교알리미 사이트 접속 필요, 데이터 시각화 부족", "building.2"),
]

for i, (title, desc, _) in enumerate(problems):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.8 + row * 2.2)
    card = add_shape(slide, x, y, Inches(3.7), Inches(1.9), BG_CARD)
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.3), Inches(0.35),
             title, 16, RED, True)
    add_text(slide, x + Inches(0.2), y + Inches(0.65), Inches(3.3), Inches(1.0),
             desc, 12, GRAY)

add_text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
         '"매일 반복되는 학교생활 정보를, 하나의 앱에서 AI와 함께 해결할 수는 없을까?"', 16, ACCENT, True, PP_ALIGN.CENTER)

# ══════════════════════════════════════
# 4. 서비스 개요
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
         "03  서비스 개요", 28, WHITE, True)

tf = add_text(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(1),
              "오늘시간표란?", 22, WHITE, True)
add_para(tf, "교육 공공데이터(NEIS/1365/학교알리미)와 AI를 활용한", 14, GRAY)
add_para(tf, "학교생활 통합 도우미 앱", 14, ACCENT, True)

features_overview = [
    ("60+ 기능", "시간표/급식/봉사/수행/AI분석/뉴스"),
    ("7종 공공데이터", "NEIS 6종 + 1365 + 학교알리미"),
    ("AI 엔진 3종", "Gemini + Groq + Claude Code"),
    ("5개 플랫폼", "iOS/iPad/macOS/Watch/위젯"),
]

for i, (title, desc) in enumerate(features_overview):
    x = Inches(0.8 + i * 3.1)
    y = Inches(3.0)
    card = add_shape(slide, x, y, Inches(2.8), Inches(1.5), BG_CARD)
    add_text(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.5), Inches(0.4),
             title, 20, ACCENT, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.5), Inches(0.6),
             desc, 12, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
         "핵심 가치:  교육 공공데이터 + AI + 멀티플랫폼 = 학생 중심의 스마트 학교생활", 14, ACCENT2, True, PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1),
         "(스크린샷: 홈 대시보드 - 인사말 + 시간표 요약 + 급식 + D-Day + 수행평가 + 날씨)", 13, GRAY, False, PP_ALIGN.CENTER)

# ══════════════════════════════════════
# 5. 핵심 기능 1 - 시간표
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
         "04  핵심 기능 ① 스마트 시간표", 28, WHITE, True)

features_tt = [
    "NEIS API로 초/중/고 시간표 자동 조회",
    "일간/주간 시간표, 주별 탐색 (스와이프)",
    "사진 촬영 → Gemini Vision AI 시간표 자동 인식",
    "과목 수정 (꾹 눌러 편집 + Firebase 동기화)",
    "과거 시간표 열람 (1개월, Firestore 이력)",
    "배경화면 자동 생성 (시간표 이미지 → 사진 앨범 저장)",
    "카카오톡/인스타그램 시간표 공유",
    "수업 시작 전 푸시 알림",
]

card = add_shape(slide, Inches(0.8), Inches(1.2), Inches(6.5), Inches(5.5), BG_CARD)
for i, f in enumerate(features_tt):
    add_text(slide, Inches(1.0), Inches(1.4 + i * 0.6), Inches(6.0), Inches(0.5),
             f"  {f}", 13, LIGHT)

# AI 시간표 인식 박스
card2 = add_shape(slide, Inches(7.8), Inches(1.2), Inches(4.7), Inches(2.5), BG_CARD)
add_text(slide, Inches(8.0), Inches(1.3), Inches(4.3), Inches(0.4),
         "AI 시간표 사진 인식", 16, ORANGE, True)
algo_text = "1. 시간표 사진 촬영/선택\n2. Gemini Vision AI 분석\n3. 과목/교시 자동 추출\n4. 시간표에 자동 입력\n\n→ 수동 입력 불필요!"
add_text(slide, Inches(8.0), Inches(1.8), Inches(4.3), Inches(2.0),
         algo_text, 11, GRAY)

add_text(slide, Inches(7.8), Inches(4.0), Inches(4.7), Inches(2.5),
         "(스크린샷: 일간 시간표\n+ 주간 시간표 격자)", 13, GRAY, False, PP_ALIGN.CENTER)

# ══════════════════════════════════════
# 6. 핵심 기능 2 - 급식 + 학사일정
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "05  핵심 기능 ② 급식 + 학사일정", 28, WHITE, True)

# 급식
card = add_shape(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.5), BG_CARD)
add_text(slide, Inches(1.0), Inches(1.3), Inches(5.4), Inches(0.4),
         "급식 (NEIS API)", 18, ACCENT2, True)
meal_features = [
    "조식/중식/석식 + 칼로리 표시",
    "19종 알레르기 빨간색 강조",
    "AI 급식 이미지 생성 (Gemini)",
    "AI 운동 분석 (칼로리 소비 추천)",
    "HealthKit 연동 (건강앱 칼로리 기록)",
    "카카오톡/인스타 공유",
]
for i, f in enumerate(meal_features):
    add_text(slide, Inches(1.0), Inches(1.9 + i * 0.5), Inches(5.4), Inches(0.4),
             f"  {f}", 13, LIGHT)

# 학사일정
card2 = add_shape(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.5), BG_CARD)
add_text(slide, Inches(7.2), Inches(1.3), Inches(5.4), Inches(0.4),
         "학사일정 + D-Day (NEIS API)", 18, ORANGE, True)
cal_features = [
    "월간 캘린더 + 스와이프 전환",
    "휴업일/체험학습/시험 분류 표시",
    "시험 D-Day 자동 감지 (키워드)",
    "iOS 캘린더 자동 추가 (알림 포함)",
    "D-Day 인스타 스토리 공유",
    "내일 학사일정 푸시 알림",
]
for i, f in enumerate(cal_features):
    add_text(slide, Inches(7.2), Inches(1.9 + i * 0.5), Inches(5.4), Inches(0.4),
             f"  {f}", 13, LIGHT)

# ══════════════════════════════════════
# 7. 핵심 기능 3 - 봉사 + 수행평가
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "06  핵심 기능 ③ 봉사활동 + 수행평가 AI", 28, WHITE, True)

# 봉사
card = add_shape(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.5), BG_CARD)
add_text(slide, Inches(1.0), Inches(1.3), Inches(5.4), Inches(0.4),
         "봉사활동 (1365 API)", 18, GREEN, True)
vol_features = [
    "지역별 봉사 프로그램 검색",
    "모집 상태/일정/기관 정보 확인",
    "Firebase 프록시 (HTTP→HTTPS, ATS 대응)",
    "6시간 캐시로 API 부하 최소화",
    "AI 봉사활동 추천 (관심 분야 기반)",
]
for i, f in enumerate(vol_features):
    add_text(slide, Inches(1.0), Inches(1.9 + i * 0.55), Inches(5.4), Inches(0.4),
             f"  {f}", 13, LIGHT)

# 수행평가
card2 = add_shape(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.5), BG_CARD)
add_text(slide, Inches(7.2), Inches(1.3), Inches(5.4), Inches(0.4),
         "수행평가 AI 도우미 (Gemini Vision)", 18, PURPLE, True)
perf_features = [
    "안내문 사진 촬영 → AI 자동 분석",
    "과목/날짜/평가기준/준비물 자동 추출",
    "시간표에서 교시 자동 매칭",
    "AI 준비 팁 + AI 코칭 (상세 전략)",
    "D-Day 알림 (D-7, D-3, D-1, 당일)",
    "Gemini 429 시 Groq 자동 전환",
]
for i, f in enumerate(perf_features):
    add_text(slide, Inches(7.2), Inches(1.9 + i * 0.55), Inches(5.4), Inches(0.4),
             f"  {f}", 13, LIGHT)

# ══════════════════════════════════════
# 8. AI 분석 기능 5종
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "07  AI 분석 기능 (5종)", 28, WHITE, True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
         "교육 공공데이터 + AI = 학생 맞춤형 분석", 14, GRAY)

ai_features = [
    ("AI 공부 플래너", "시험 정보 입력 → 일별 과목별\n공부 계획표 자동 생성\n(Swift Charts 시각화)", ACCENT),
    ("AI 영양 리포트", "주간 급식 데이터 분석 →\n영양 균형/부족 영양소/\n보충 음식 추천", ACCENT2),
    ("AI 봉사 추천", "관심 분야 기반\n봉사활동 추천 + 1365\n검색 키워드 제안", GREEN),
    ("AI 학습 분석", "집중 모드 기록 분석 →\n학습 패턴/개선점/\n목표 추천 (7일 차트)", PURPLE),
    ("AI 학교 비교", "NEIS 검색으로 학교 선택 →\n학교알리미 데이터 기반\n비교 차트 + AI 분석", ORANGE),
]

for i, (title, desc, color) in enumerate(ai_features):
    x = Inches(0.5 + i * 2.5)
    card = add_shape(slide, x, Inches(1.6), Inches(2.3), Inches(3.5), BG_CARD)
    # 상단 색 바
    bar = add_shape(slide, x, Inches(1.6), Inches(2.3), Inches(0.06), color)
    add_text(slide, x + Inches(0.1), Inches(1.8), Inches(2.1), Inches(0.4),
             title, 15, color, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(2.4), Inches(2.1), Inches(2.5),
             desc, 11, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4),
         "AI 엔진: Gemini Flash Lite (우선) → Groq Llama 3.3 70B (fallback, 한국어 시스템 프롬프트)", 12, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.8),
         "(스크린샷: AI 공부 플래너 차트 + AI 영양 리포트 차트 + AI 학교 비교 차트)", 13, GRAY, False, PP_ALIGN.CENTER)

# ══════════════════════════════════════
# 9. 학교알리미 + 교육 뉴스
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "08  학교알리미 정보공시 + 교육 뉴스", 28, WHITE, True)

# 학교알리미
card = add_shape(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.5), BG_CARD)
add_text(slide, Inches(1.0), Inches(1.3), Inches(5.4), Inches(0.4),
         "학교알리미 정보공시 (12종 API)", 17, ACCENT, True)
school_info = [
    "학생 현황 (학년별 반수/평균 학생수)",
    "성별 현황 (남녀 비율 바 그래프)",
    "교원 현황 (교장/교감/교사 수)",
    "동아리 목록 + 인원수",
    "학교도서관 (장서/좌석)",
    "교복 단가",
    "방과후학교 프로그램",
    "수업일수 + 학년별 학생수 상세",
    "AI 학교 진단 (데이터 기반 AI 분석)",
    "지도 (MapKit) + 연락처",
]
for i, f in enumerate(school_info):
    add_text(slide, Inches(1.0), Inches(1.9 + i * 0.45), Inches(5.4), Inches(0.35),
             f"  {f}", 12, LIGHT)

# 뉴스
card2 = add_shape(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.5), BG_CARD)
add_text(slide, Inches(7.2), Inches(1.3), Inches(5.4), Inches(0.4),
         "교육 뉴스 (카드뉴스)", 17, ORANGE, True)
news_features = [
    "교육정책/입시/학교생활/진로 카테고리",
    "카드 스타일 뉴스 피드 + 상세 보기",
    "이미지 + 글 지원",
    "새 뉴스 로컬 알림 (사용자 동의 후)",
    "관리자 모드 (설정 > 버전 7번 탭)",
    "Firebase Firestore + Cloud Functions",
    "공공누리 출처 표시",
    "교육부 정보공시 기준 안내",
]
for i, f in enumerate(news_features):
    add_text(slide, Inches(7.2), Inches(1.9 + i * 0.5), Inches(5.4), Inches(0.4),
             f"  {f}", 13, LIGHT)

# ══════════════════════════════════════
# 10. 멀티플랫폼 + 위젯
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "09  멀티플랫폼 + 위젯 + Live Activity", 28, WHITE, True)

platforms = [
    ("iPhone", "전체 기능\n홈/시간표/급식/봉사\n수행평가/AI/뉴스", ACCENT),
    ("iPad", "NavigationSplitView\n사이드바 레이아웃\n넓은 화면 최적화", ACCENT2),
    ("Apple Watch", "시간표 + 급식\n독립 실행\n오프라인 알림", GREEN),
    ("macOS", "메뉴바 앱\n풀 윈도우 모드\nMac App Store", PURPLE),
]

for i, (name, desc, color) in enumerate(platforms):
    x = Inches(0.8 + i * 3.1)
    card = add_shape(slide, x, Inches(1.2), Inches(2.8), Inches(2.5), BG_CARD)
    add_text(slide, x + Inches(0.1), Inches(1.35), Inches(2.6), Inches(0.4),
             name, 20, color, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(1.85), Inches(2.6), Inches(1.5),
             desc, 12, GRAY, False, PP_ALIGN.CENTER)

# 위젯 + LA
widgets = [
    ("위젯 6종", "시간표/급식/학사일정/올인원\n(홈화면 + 잠금화면 + macOS)"),
    ("Live Activity", "다이나믹 아일랜드에 현재 수업\n집중 모드 타이머 표시"),
    ("기기 동기화", "iCloud Key-Value Store\n+ Firebase Firestore"),
    ("공유", "카카오톡 + 인스타 스토리\nD-Day 이미지 공유"),
]

for i, (title, desc) in enumerate(widgets):
    x = Inches(0.8 + i * 3.1)
    card = add_shape(slide, x, Inches(4.2), Inches(2.8), Inches(2.0), BG_CARD)
    add_text(slide, x + Inches(0.1), Inches(4.35), Inches(2.6), Inches(0.35),
             title, 15, ACCENT, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(4.8), Inches(2.6), Inches(1.2),
             desc, 11, GRAY, False, PP_ALIGN.CENTER)

# ══════════════════════════════════════
# 11. 집중 모드 + 건강
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "10  집중 모드 + 건강 연동 + 배경화면", 28, WHITE, True)

# 집중 모드
card = add_shape(slide, Inches(0.8), Inches(1.2), Inches(3.7), Inches(5.0), BG_CARD)
add_text(slide, Inches(1.0), Inches(1.3), Inches(3.3), Inches(0.4),
         "집중 모드", 18, PURPLE, True)
focus_items = ["과목별 공부 스톱워치", "일별/주별 공부 시간 기록", "Live Activity 타이머", "방해금지 모드 연동"]
for i, f in enumerate(focus_items):
    add_text(slide, Inches(1.0), Inches(1.9 + i * 0.5), Inches(3.3), Inches(0.4),
             f"  {f}", 12, LIGHT)

# 건강
card2 = add_shape(slide, Inches(4.8), Inches(1.2), Inches(3.7), Inches(5.0), BG_CARD)
add_text(slide, Inches(5.0), Inches(1.3), Inches(3.3), Inches(0.4),
         "건강 연동", 18, RED, True)
health_items = ["급식 칼로리 → HealthKit 기록", "활동 칼로리 확인", "AI 운동 분석 (급식 기반)", "날씨 정보 (Open-Meteo)"]
for i, f in enumerate(health_items):
    add_text(slide, Inches(5.0), Inches(1.9 + i * 0.5), Inches(3.3), Inches(0.4),
             f"  {f}", 12, LIGHT)

# 배경화면
card3 = add_shape(slide, Inches(8.8), Inches(1.2), Inches(3.7), Inches(5.0), BG_CARD)
add_text(slide, Inches(9.0), Inches(1.3), Inches(3.3), Inches(0.4),
         "배경화면 자동 생성", 18, ACCENT2, True)
wall_items = ["주간 시간표 이미지 렌더링", "사진 앨범 자동 저장", "iOS 단축어 연동 가이드", "시간표 새로고침 시 자동 생성"]
for i, f in enumerate(wall_items):
    add_text(slide, Inches(9.0), Inches(1.9 + i * 0.5), Inches(3.3), Inches(0.4),
             f"  {f}", 12, LIGHT)

# ══════════════════════════════════════
# 12. 기술 아키텍처
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "11  기술 아키텍처", 28, WHITE, True)

# 아키텍처 다이어그램 (텍스트)
arch_box = add_shape(slide, Inches(0.8), Inches(1.2), Inches(7.5), Inches(5.5), BG_CARD)

add_text(slide, Inches(1.0), Inches(1.4), Inches(7.0), Inches(0.4),
         "사용자 앱 (iOS / macOS / watchOS)", 16, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(3.5), Inches(1.9), Inches(2.5), Inches(0.3),
         "HTTPS", 12, ACCENT, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.0), Inches(2.3), Inches(7.0), Inches(0.4),
         "Firebase Cloud Functions (13개)", 16, ACCENT, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.0), Inches(2.8), Inches(7.0), Inches(0.3),
         "API 키는 서버에만 보관 (Secret Manager)", 11, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.0), Inches(3.4), Inches(3.3), Inches(0.4),
         "Firestore (캐시 DB)", 14, ACCENT2, True, PP_ALIGN.CENTER)
add_text(slide, Inches(4.5), Inches(3.4), Inches(3.3), Inches(0.4),
         "Firebase Storage (이미지)", 14, ACCENT2, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1.0), Inches(4.2), Inches(7.0), Inches(0.3),
         "외부 API", 14, ORANGE, True, PP_ALIGN.CENTER)

apis = ["NEIS API", "학교알리미", "1365 봉사", "Gemini AI", "Groq AI", "Open-Meteo"]
for i, api in enumerate(apis):
    x = Inches(1.0 + i * 1.15)
    add_badge(slide, x, Inches(4.7), api, ORANGE if i < 3 else PURPLE)

# Cloud Functions 목록
funcs_box = add_shape(slide, Inches(8.8), Inches(1.2), Inches(4.0), Inches(5.5), BG_CARD)
add_text(slide, Inches(9.0), Inches(1.3), Inches(3.6), Inches(0.4),
         "Cloud Functions 13개", 15, ACCENT, True)
funcs = [
    "1. 학교 검색", "2. NEIS 시간표", "3. 급식 조회",
    "4. 학사일정", "5. 시간표 이력 저장", "6. 시간표 수정 동기화",
    "7. 시간표 이력", "8-9. 수정 동기화",
    "10-11. 봉사활동", "12. 뉴스 조회", "13. 뉴스 발행",
]
for i, f in enumerate(funcs):
    add_text(slide, Inches(9.0), Inches(1.8 + i * 0.42), Inches(3.6), Inches(0.35),
             f, 11, GRAY)

# ══════════════════════════════════════
# 13. AI 활용 상세
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "12  AI 활용 상세", 28, WHITE, True)

ai_tools = [
    ("Claude Code\n(Opus 4.6)", "앱 전체 개발", "Swift/TypeScript/Kotlin\n코드 생성, 디버깅\n아키텍처 설계\n100+ 회 대화 협업", ACCENT),
    ("Gemini\nFlash Lite", "텍스트 AI 분석", "수행평가 분석/AI 코칭\n공부 플래너/영양 리포트\n학교 비교/학습 분석\n3개 키 로테이션", ACCENT2),
    ("Gemini\nFlash Image", "이미지 생성", "급식 메뉴 → 식판 이미지\n한국어 프롬프트\n\"한국 학교 급식 식판\"", ORANGE),
    ("Gemini\nVision", "이미지 인식", "시간표 사진 → 데이터 변환\n수행평가 안내문 → JSON\n과목/날짜/기준 자동 추출", PURPLE),
    ("Groq\nLlama 3.3 70B", "텍스트 Fallback", "Gemini 429 시 자동 전환\n한국어 시스템 프롬프트\n무료 + 빠른 응답", GREEN),
]

for i, (name, role, desc, color) in enumerate(ai_tools):
    x = Inches(0.4 + i * 2.55)
    card = add_shape(slide, x, Inches(1.2), Inches(2.35), Inches(4.5), BG_CARD)
    bar = add_shape(slide, x, Inches(1.2), Inches(2.35), Inches(0.06), color)
    add_text(slide, x + Inches(0.1), Inches(1.4), Inches(2.15), Inches(0.7),
             name, 14, color, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(2.1), Inches(2.15), Inches(0.35),
             role, 11, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(2.6), Inches(2.15), Inches(2.8),
             desc, 10, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
         "AI 활용 비율: 코드의 약 90%를 AI와 협업하여 작성  |  보안: API 키 3중 로테이션 + Groq fallback", 12, GRAY, False, PP_ALIGN.CENTER)

# ══════════════════════════════════════
# 14. 차별점 + 기대효과
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         "13  기존 서비스 차별점 + 기대효과", 28, WHITE, True)

# 비교 표
headers = ["비교 항목", "기존 시간표 앱", "오늘시간표"]
cols_x = [Inches(1.5), Inches(5.0), Inches(8.5)]
cols_w = [Inches(3.2), Inches(3.2), Inches(3.5)]

for i, (h, x, w) in enumerate(zip(headers, cols_x, cols_w)):
    add_text(slide, x, Inches(1.1), w, Inches(0.35), h, 13, ACCENT if i == 2 else GRAY, True, PP_ALIGN.CENTER)

rows = [
    ("시간표 조회", "일부", "O (NEIS 초/중/고)"),
    ("급식 + AI 이미지", "일부", "O (Gemini 이미지+운동분석)"),
    ("수행평가 AI", "X", "O (Vision+코칭)"),
    ("봉사활동", "X", "O (1365 API)"),
    ("AI 분석 5종", "X", "O (차트 포함)"),
    ("학교알리미 12종", "X", "O + AI 진단"),
    ("교육 뉴스", "X", "O + 알림"),
    ("Apple Watch", "X", "O (독립 실행)"),
    ("위젯 6종 + Live Activity", "일부", "O"),
]

for j, (item, c1, c2) in enumerate(rows):
    y = Inches(1.5 + j * 0.42)
    for x, w, val in [(cols_x[0], cols_w[0], item), (cols_x[1], cols_w[1], c1),
                       (cols_x[2], cols_w[2], c2)]:
        color = GREEN if val.startswith("O") else (RED if val == "X" else LIGHT)
        add_text(slide, x, y, w, Inches(0.35), val, 11, color, val.startswith("O"), PP_ALIGN.CENTER)

# 기대효과
add_text(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4),
         "기대효과", 18, WHITE, True)
effects = [
    "학생: 정보 확인 시간 절약 + AI 자기주도 학습 + 봉사활동 관리",
    "교사: \"오늘 시간표 뭐예요?\" 질문 감소 + 변경 공지 부담 경감",
    "사회: 교육 공공데이터의 실질적 활용 사례 + 전국 6,000+ 학교 확장 가능",
]
for i, e in enumerate(effects):
    add_text(slide, Inches(0.8), Inches(6.0 + i * 0.4), Inches(11.5), Inches(0.35),
             f"  {e}", 12, LIGHT)

# ══════════════════════════════════════
# 15. 요약
# ══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8),
         "오늘시간표", 44, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(1.7), Inches(10), Inches(0.5),
         "교육 공공데이터 + AI + 학생 중심 설계", 20, ACCENT, False, PP_ALIGN.CENTER)

summary = [
    ("활용 데이터", "NEIS 6종 + 1365 봉사 + 학교알리미 12종"),
    ("AI 기술", "Claude Code (개발) + Gemini (분석/이미지) + Groq (fallback)"),
    ("핵심 기능", "시간표(변경감지) + 급식(AI) + 봉사 + 수행평가AI + AI분석 5종 + 뉴스"),
    ("플랫폼", "iOS + iPad + macOS + watchOS + 위젯 6종 + Live Activity"),
    ("기능 수", "60+ 기능 구현 완료"),
]

for i, (label, value) in enumerate(summary):
    y = Inches(2.5 + i * 0.65)
    card = add_shape(slide, Inches(1.5), y, Inches(10), Inches(0.55), BG_CARD)
    add_text(slide, Inches(1.7), y + Inches(0.08), Inches(2.5), Inches(0.4),
             label, 14, ACCENT, True)
    add_text(slide, Inches(4.2), y + Inches(0.08), Inches(7.0), Inches(0.4),
             value, 13, LIGHT)

add_text(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
         '"학생이 만든, 학생을 위한, AI 기반 학교생활 도우미"', 18, ACCENT2, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
         "(App Store QR코드 + 다운로드 링크)", 14, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
         "공공누리 제1유형  |  출처: NEIS(open.neis.go.kr), 1365(1365.go.kr), 학교알리미(schoolinfo.go.kr)", 9, GRAY, False, PP_ALIGN.CENTER)

# === 저장 ===
output = "/Users/sin-yeseong/School/TodayTimetable/대회/오늘시간표_기획서_v2.pptx"
prs.save(output)
print(f"PPT 생성 완료: {output}")
print(f"총 {len(prs.slides)}매")
